"""
Generates LaerdalImplementation.pptx — a presentation covering the architecture,
design decisions, and key flows of the Laerdal Implementer backend.
Open in Keynote (File > Open) and save as .key if needed.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Colour palette ────────────────────────────────────────────────────────────
LAERDAL_RED   = RGBColor(0xC0, 0x20, 0x2C)   # Laerdal brand red
DARK_BG       = RGBColor(0x1A, 0x1A, 0x2E)   # deep navy
MID_BG        = RGBColor(0x16, 0x21, 0x3E)   # section bg
ACCENT        = RGBColor(0x0F, 0x3A, 0x60)   # card bg
WHITE         = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY    = RGBColor(0xD0, 0xD5, 0xDD)
YELLOW        = RGBColor(0xFF, 0xC1, 0x07)
GREEN         = RGBColor(0x28, 0xA7, 0x45)
BLUE_ACCENT   = RGBColor(0x17, 0xA2, 0xB8)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

blank_layout = prs.slide_layouts[6]   # completely blank


# ── Helper functions ──────────────────────────────────────────────────────────

def add_slide():
    return prs.slides.add_slide(blank_layout)

def bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def box(slide, x, y, w, h, fill_color=None, line_color=None, line_width=Pt(0)):
    shape = slide.shapes.add_shape(1, x, y, w, h)   # MSO_SHAPE_TYPE.RECTANGLE = 1
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape

def txt(slide, text, x, y, w, h,
        size=Pt(18), bold=False, color=WHITE, align=PP_ALIGN.LEFT,
        italic=False, wrap=True):
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = size
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txb

def divider(slide, y, color=LAERDAL_RED, thickness=Pt(2)):
    line = slide.shapes.add_shape(1, Inches(0.5), y, SLIDE_W - Inches(1), thickness)
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()

def bullet_block(slide, items, x, y, w, h,
                 size=Pt(16), color=LIGHT_GREY, bullet="•  ", bold_first=False):
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.space_before = Pt(4)
        run = p.add_run()
        if isinstance(item, tuple):
            label, rest = item
            run.text = bullet + label
            run.font.bold  = True
            run.font.color.rgb = WHITE
            run.font.size  = size
            if rest:
                run2 = p.add_run()
                run2.text = rest
                run2.font.bold  = False
                run2.font.color.rgb = color
                run2.font.size  = size
        else:
            run.text = bullet + item
            run.font.bold  = bold_first and i == 0
            run.font.color.rgb = color
            run.font.size  = size
    return txb

def card(slide, x, y, w, h, title, body_lines,
         title_color=WHITE, title_size=Pt(17),
         body_size=Pt(14), accent_color=LAERDAL_RED):
    box(slide, x, y, w, h, fill_color=ACCENT,
        line_color=accent_color, line_width=Pt(1.5))
    # title
    txt(slide, title, x + Inches(0.15), y + Inches(0.12), w - Inches(0.3), Inches(0.4),
        size=title_size, bold=True, color=title_color)
    # divider rule inside card
    rule = slide.shapes.add_shape(1, x + Inches(0.15), y + Inches(0.5),
                                   w - Inches(0.3), Pt(1))
    rule.fill.solid(); rule.fill.fore_color.rgb = accent_color
    rule.line.fill.background()
    # body
    bullet_block(slide, body_lines,
                 x + Inches(0.15), y + Inches(0.58),
                 w - Inches(0.3), h - Inches(0.7),
                 size=body_size, bullet="")


def section_header(slide, number, title, subtitle=""):
    bg(slide, MID_BG)
    # big number watermark
    txt(slide, number, Inches(0.3), Inches(0.5), Inches(2), Inches(5),
        size=Pt(200), bold=True, color=RGBColor(0x25, 0x35, 0x55), align=PP_ALIGN.LEFT)
    txt(slide, title, Inches(2.8), Inches(2.4), Inches(9.5), Inches(1.2),
        size=Pt(48), bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        txt(slide, subtitle, Inches(2.8), Inches(3.5), Inches(9.5), Inches(0.7),
            size=Pt(22), bold=False, color=LIGHT_GREY, align=PP_ALIGN.LEFT)
    divider(slide, Inches(2.2), color=LAERDAL_RED, thickness=Pt(4))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s, DARK_BG)

# red accent bar left
box(s, Inches(0), Inches(0), Inches(0.25), SLIDE_H, fill_color=LAERDAL_RED)

txt(s, "LAERDAL IMPLEMENTER", Inches(0.6), Inches(1.6), Inches(12), Inches(0.9),
    size=Pt(42), bold=True, color=LAERDAL_RED, align=PP_ALIGN.LEFT)
txt(s, "Backend Service Design", Inches(0.6), Inches(2.45), Inches(12), Inches(0.9),
    size=Pt(52), bold=True, color=WHITE, align=PP_ALIGN.LEFT)
divider(s, Inches(3.35), color=LAERDAL_RED, thickness=Pt(3))
txt(s, "Organization & Course Management · C# / .NET 8",
    Inches(0.6), Inches(3.6), Inches(12), Inches(0.5),
    size=Pt(22), color=LIGHT_GREY, align=PP_ALIGN.LEFT)

# three mini-pills at bottom
pill_labels = ["Clean Architecture", "CQRS + MediatR", "EF Core + SQL Server"]
for i, label in enumerate(pill_labels):
    px = Inches(0.6) + i * Inches(3.5)
    box(s, px, Inches(6.4), Inches(3.1), Inches(0.7),
        fill_color=ACCENT, line_color=LAERDAL_RED, line_width=Pt(1))
    txt(s, label, px + Inches(0.1), Inches(6.45), Inches(3), Inches(0.55),
        size=Pt(16), bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — What does this system do?
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s, DARK_BG)
box(s, Inches(0), Inches(0), Inches(0.25), SLIDE_H, fill_color=LAERDAL_RED)

txt(s, "What does this system do?", Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
    size=Pt(36), bold=True, color=WHITE)
divider(s, Inches(1.0), color=LAERDAL_RED)

cards_data = [
    ("1  Manage Organizations",
     ["Create, list, update, soft-delete hospitals & customers",
      "Parent → child hierarchy (Hospital → Department)",
      "Each org has a unique Code within its parent scope",
      "Soft-delete preserves audit trail & FK integrity"]),
    ("2  Manage Manifests",
     ["Versioned bundles of course + learning-activity content",
      "Lifecycle: Draft → Published → Archived",
      "Published manifests are immutable snapshots",
      "Training app pulls the active manifest per org"]),
    ("3  Authentication",
     ["React frontend authenticates via Laerdal OIDC (JWT)",
      "Role-based: Laerdal Admin vs Hospital Admin",
      "Training app uses OAuth 2.0 Client Credentials (S2S)",
      "Policy-based authorization in ASP.NET Core"]),
]

for i, (title, lines) in enumerate(cards_data):
    cx = Inches(0.5) + i * Inches(4.2)
    card(s, cx, Inches(1.25), Inches(3.9), Inches(5.8), title, lines,
         title_size=Pt(16), body_size=Pt(14))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Section: Architecture
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
section_header(s, "01", "Clean Architecture", "Why layers? What's in each one?")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Architecture diagram
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s, DARK_BG)
box(s, Inches(0), Inches(0), Inches(0.25), SLIDE_H, fill_color=LAERDAL_RED)

txt(s, "Four-Layer Architecture", Inches(0.6), Inches(0.2), Inches(10), Inches(0.6),
    size=Pt(32), bold=True, color=WHITE)
divider(s, Inches(0.85), color=LAERDAL_RED)

txt(s, "The dependency rule: outer layers know about inner layers — never the reverse.",
    Inches(0.6), Inches(0.9), Inches(12), Inches(0.45),
    size=Pt(16), color=LIGHT_GREY, italic=True)

# Layer boxes — drawn right to left (outermost first) so inner ones sit on top visually
layers = [
    # (label, sublabel, color, x)
    ("API",            "Controllers · Models · Program.cs",                     RGBColor(0xC0,0x20,0x2C), Inches(0.55)),
    ("Infrastructure", "DbContext · Repositories · Migrations · Config",        RGBColor(0x0F,0x3A,0x60), Inches(3.45)),
    ("Application",    "Commands · Queries · DTOs · Mappers",                   RGBColor(0x1A,0x52,0x76), Inches(6.35)),
    ("Domain",         "Entities · Enums · Repository Interfaces",              RGBColor(0x28,0x7A,0x56), Inches(9.25)),
]

for label, sub, color, x in layers:
    box(s, x, Inches(1.6), Inches(2.6), Inches(4.2),
        fill_color=color, line_color=WHITE, line_width=Pt(0.5))
    txt(s, label, x + Inches(0.1), Inches(1.75), Inches(2.4), Inches(0.6),
        size=Pt(22), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, sub, x + Inches(0.1), Inches(2.35), Inches(2.4), Inches(1.0),
        size=Pt(12), color=RGBColor(0xE0,0xE8,0xF0), align=PP_ALIGN.CENTER, italic=True)

# arrows between layers
for ax in [Inches(3.1), Inches(6.0), Inches(8.9)]:
    txt(s, "→", ax, Inches(2.8), Inches(0.5), Inches(0.5),
        size=Pt(28), bold=True, color=LIGHT_GREY, align=PP_ALIGN.CENTER)

# dependency rule callout boxes
callouts = [
    (Inches(0.55),  Inches(6.1), "HTTP layer only", "Talks to MediatR, returns ActionResult"),
    (Inches(3.45),  Inches(6.1), "SQL Server detail", "Only place that knows about EF Core"),
    (Inches(6.35),  Inches(6.1), "Use-case logic",   "Orchestrates domain + repos via interfaces"),
    (Inches(9.25),  Inches(6.1), "No dependencies",  "Pure C#, zero framework references"),
]
for (cx, cy, title, desc) in callouts:
    box(s, cx, cy, Inches(2.6), Inches(1.1),
        fill_color=ACCENT, line_color=YELLOW, line_width=Pt(1))
    txt(s, title, cx + Inches(0.1), cy + Inches(0.05), Inches(2.4), Inches(0.35),
        size=Pt(12), bold=True, color=YELLOW)
    txt(s, desc,  cx + Inches(0.1), cy + Inches(0.38), Inches(2.4), Inches(0.65),
        size=Pt(11), color=LIGHT_GREY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Section: Domain Model
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
section_header(s, "02", "Domain Model", "Organizations · Manifests · Business Rules")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Organization entity
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s, DARK_BG)
box(s, Inches(0), Inches(0), Inches(0.25), SLIDE_H, fill_color=LAERDAL_RED)

txt(s, "Organization Entity", Inches(0.6), Inches(0.2), Inches(10), Inches(0.6),
    size=Pt(32), bold=True, color=WHITE)
divider(s, Inches(0.85), color=LAERDAL_RED)

# Fields card
card(s, Inches(0.5), Inches(1.05), Inches(4.0), Inches(5.8), "Fields", [
    "Id  — GUID (generated in C#)",
    "Name  — display name",
    "Code  — short ID, stored UPPERCASE",
    "Type  — Hospital | Department | TrainingCenter",
    "ParentId  — null = root org",
    "IsActive  — soft-delete flag",
    "ExternalId  — link to OIDC/HR system",
    "CreatedAt / UpdatedAt  — audit trail",
])

# Hierarchy card
card(s, Inches(4.7), Inches(1.05), Inches(3.9), Inches(2.75), "Hierarchy (Adjacency List)", [
    "City Hospital  (ParentId = null)",
    "  └─ ICU Dept  (ParentId = hospital.Id)",
    "  └─ Training Centre  (ParentId = hospital.Id)",
    "",
    "One table. Recursive self-FK.",
    "Children loaded via EF Include().",
])

# Business rules card
card(s, Inches(4.7), Inches(3.95), Inches(3.9), Inches(2.9), "Business Rules", [
    "Code unique within parent scope",
    "CanBeDeleted() → no active children",
    "Soft-delete keeps FK references intact",
    "Factory method Create() validates on construction",
], accent_color=YELLOW)

# code snippet card
card(s, Inches(8.75), Inches(1.05), Inches(4.3), Inches(5.8), "Factory Method (Domain)", [
    "public static Organization Create(",
    "  string name, string code,",
    "  OrganizationType type,",
    "  Guid? parentId = null,",
    "  string? externalId = null)",
    "{",
    "  if (string.IsNullOrWhiteSpace(name))",
    "    throw new ArgumentException(...);",
    "",
    "  return new Organization {",
    "    Id = Guid.NewGuid(),",
    "    Code = code.Trim().ToUpper(),",
    "    IsActive = true,",
    "    CreatedAt = DateTime.UtcNow",
    "  };",
    "}",
], accent_color=BLUE_ACCENT, body_size=Pt(11))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Manifest entity
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s, DARK_BG)
box(s, Inches(0), Inches(0), Inches(0.25), SLIDE_H, fill_color=LAERDAL_RED)

txt(s, "Manifest Entity", Inches(0.6), Inches(0.2), Inches(10), Inches(0.6),
    size=Pt(32), bold=True, color=WHITE)
divider(s, Inches(0.85), color=LAERDAL_RED)

# Lifecycle diagram
txt(s, "Lifecycle:", Inches(0.6), Inches(1.05), Inches(2), Inches(0.45),
    size=Pt(16), bold=True, color=LIGHT_GREY)

stages = [("Draft", ACCENT, Inches(2.5)), ("Published", RGBColor(0x15,0x60,0x2B), Inches(5.5)), ("Archived", RGBColor(0x50,0x30,0x10), Inches(8.5))]
for label, color, lx in stages:
    box(s, lx, Inches(0.95), Inches(2.5), Inches(0.7),
        fill_color=color, line_color=LAERDAL_RED, line_width=Pt(1.5))
    txt(s, label, lx, Inches(0.95), Inches(2.5), Inches(0.7),
        size=Pt(20), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

for ax in [Inches(4.85), Inches(7.85)]:
    txt(s, "→", ax, Inches(0.95), Inches(0.65), Inches(0.7),
        size=Pt(26), bold=True, color=YELLOW, align=PP_ALIGN.CENTER)

# Fields card
card(s, Inches(0.5), Inches(1.95), Inches(4.0), Inches(5.0), "Fields", [
    "Id  — GUID",
    "OrganizationId  — FK to owner",
    "Name / Description",
    "Version  — semantic (e.g. 1.2.0)",
    "Status  — Draft | Published | Archived",
    "Content  — JSON blob (courses + activities)",
    "PublishedAt  — set once, never changed",
    "CreatedAt / UpdatedAt",
])

# Immutability card
card(s, Inches(4.65), Inches(1.95), Inches(4.0), Inches(2.4), "Immutability Rule", [
    "Published manifests CANNOT be edited.",
    "Any change → new manifest + new version.",
    "Old Published → Archive().",
    "Learner sessions pin a ManifestId.",
    "Session content never shifts mid-way.",
], accent_color=YELLOW)

# Content structure
card(s, Inches(4.65), Inches(4.5), Inches(4.0), Inches(2.45), "Content JSON Structure", [
    '{ "courses": [{',
    '    "id": "uuid",',
    '    "title": "Course Title",',
    '    "activities": [{',
    '      "type": "simulation",',
    '      "duration": 30 }]',
    '}]}',
], accent_color=BLUE_ACCENT, body_size=Pt(12))

# Why JSON blob card
card(s, Inches(8.8), Inches(1.95), Inches(4.2), Inches(5.0), "Why a JSON Blob?", [
    "Training app consumes the whole manifest",
    "in one request — no nested joins needed.",
    "",
    "Alternative: separate Course + Activity tables",
    "→ more complex queries, more migrations,",
    "  no benefit for this read pattern.",
    "",
    "Trade-off accepted: lose SQL-level querying",
    "inside content, gain simplicity.",
], accent_color=BLUE_ACCENT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Section: CQRS
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
section_header(s, "03", "CQRS + MediatR", "Separating reads from writes")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — CQRS Pattern
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s, DARK_BG)
box(s, Inches(0), Inches(0), Inches(0.25), SLIDE_H, fill_color=LAERDAL_RED)

txt(s, "CQRS with MediatR", Inches(0.6), Inches(0.2), Inches(10), Inches(0.6),
    size=Pt(32), bold=True, color=WHITE)
divider(s, Inches(0.85), color=LAERDAL_RED)

txt(s, "Command Query Responsibility Segregation — every operation is either a write (Command) or a read (Query).",
    Inches(0.6), Inches(0.92), Inches(12.3), Inches(0.45),
    size=Pt(15), color=LIGHT_GREY, italic=True)

# Commands side
box(s, Inches(0.5), Inches(1.55), Inches(5.8), Inches(5.3),
    fill_color=RGBColor(0x28,0x10,0x10), line_color=LAERDAL_RED, line_width=Pt(1.5))
txt(s, "Commands  (Writes)", Inches(0.65), Inches(1.65), Inches(5.5), Inches(0.5),
    size=Pt(20), bold=True, color=LAERDAL_RED)

cmd_items = [
    "CreateOrganizationCommand",
    "  → validates uniqueness + parent",
    "  → calls Organization.Create()",
    "  → persists via IOrganizationRepository",
    "",
    "PublishManifestCommand  (stub)",
    "  → calls manifest.Publish(version)",
    "  → archives previous Published manifest",
    "  → saves both in one transaction",
]
bullet_block(s, cmd_items, Inches(0.7), Inches(2.2), Inches(5.4), Inches(4.2),
             size=Pt(14), color=LIGHT_GREY, bullet="")

# Queries side
box(s, Inches(6.7), Inches(1.55), Inches(6.1), Inches(5.3),
    fill_color=RGBColor(0x08,0x28,0x18), line_color=GREEN, line_width=Pt(1.5))
txt(s, "Queries  (Reads)", Inches(6.85), Inches(1.65), Inches(5.8), Inches(0.5),
    size=Pt(20), bold=True, color=GREEN)

qry_items = [
    "GetOrganizationsQuery",
    "  → optional ?parentId filter",
    "  → returns full subtree via Include()",
    "",
    "GetOrganizationByIdQuery",
    "  → single lookup by PK",
    "  → uses GetByIdAsync (not a full scan)",
    "",
    "GetPublishedManifestQuery  (stub)",
    "  → used by training app",
    "  → hits IX_Manifest_OrgStatus index",
]
bullet_block(s, qry_items, Inches(6.9), Inches(2.2), Inches(5.7), Inches(4.2),
             size=Pt(14), color=LIGHT_GREY, bullet="")

# MediatR label in middle
txt(s, "MediatR\ndispatches", Inches(5.6), Inches(3.1), Inches(1.2), Inches(0.9),
    size=Pt(12), bold=True, color=YELLOW, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Section: End-to-End Flow
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
section_header(s, "04", "Create Organisation — End to End", "Tracing one request through all four layers")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Flow diagram
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s, DARK_BG)
box(s, Inches(0), Inches(0), Inches(0.25), SLIDE_H, fill_color=LAERDAL_RED)

txt(s, "POST /api/organizations  — Request Flow", Inches(0.6), Inches(0.2), Inches(12), Inches(0.6),
    size=Pt(30), bold=True, color=WHITE)
divider(s, Inches(0.82), color=LAERDAL_RED)

steps = [
    ("HTTP", LAERDAL_RED,
     "POST /api/organizations",
     'Body: { "name":"City Hospital", "code":"CH001", "type":0 }'),
    ("Controller", RGBColor(0x0F,0x3A,0x60),
     "OrganizationsController.Create()",
     "Builds CreateOrganizationCommand  →  _mediator.Send(command)"),
    ("Handler", RGBColor(0x1A,0x52,0x76),
     "CreateOrganizationCommandHandler.Handle()",
     "① Check code uniqueness (GetByCodeAsync)  ② Check parent exists (ExistsAsync)"),
    ("Domain", RGBColor(0x28,0x7A,0x56),
     "Organization.Create()",
     "Validates name/code  →  returns new entity with Guid.NewGuid()"),
    ("Repository", RGBColor(0x0F,0x3A,0x60),
     "OrganizationRepository.AddAsync()",
     "context.Organizations.Add(org)  →  SaveChangesAsync()  →  INSERT SQL"),
    ("Response", LAERDAL_RED,
     "201 Created",
     "Location: /api/organizations/{id}  +  OrganizationResponse JSON body"),
]

for i, (layer, color, title, desc) in enumerate(steps):
    sy = Inches(1.0) + i * Inches(1.02)
    box(s, Inches(0.5), sy, Inches(1.4), Inches(0.82),
        fill_color=color, line_color=color, line_width=Pt(0))
    txt(s, layer, Inches(0.5), sy, Inches(1.4), Inches(0.82),
        size=Pt(13), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    box(s, Inches(2.0), sy, Inches(10.8), Inches(0.82),
        fill_color=ACCENT, line_color=color, line_width=Pt(1))
    txt(s, title, Inches(2.15), sy + Inches(0.03), Inches(10.5), Inches(0.35),
        size=Pt(15), bold=True, color=WHITE)
    txt(s, desc, Inches(2.15), sy + Inches(0.38), Inches(10.5), Inches(0.38),
        size=Pt(12), color=LIGHT_GREY)
    if i < len(steps) - 1:
        txt(s, "↓", Inches(0.8), sy + Inches(0.82), Inches(0.9), Inches(0.22),
            size=Pt(14), color=color, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Section: API Surface
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
section_header(s, "05", "API Surface", "The 7 most important endpoints")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — API endpoints
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s, DARK_BG)
box(s, Inches(0), Inches(0), Inches(0.25), SLIDE_H, fill_color=LAERDAL_RED)

txt(s, "REST Endpoints", Inches(0.6), Inches(0.2), Inches(10), Inches(0.6),
    size=Pt(32), bold=True, color=WHITE)
divider(s, Inches(0.85), color=LAERDAL_RED)

endpoints = [
    ("POST",   "/api/organizations",                                        "Create organization  →  201 + Location header",                      LAERDAL_RED),
    ("GET",    "/api/organizations?parentId={id}",                          "List organizations  (all or children of a parent)  →  200",          GREEN),
    ("GET",    "/api/organizations/{id}",                                   "Get single org with hierarchy  →  200 / 404",                         GREEN),
    ("PATCH",  "/api/organizations/{id}",                                   "Update name / type / isActive  →  200  (stub)",                      YELLOW),
    ("POST",   "/api/organizations/{orgId}/manifests",                      "Create draft manifest  →  201  (stub)",                              LAERDAL_RED),
    ("POST",   "/api/organizations/{orgId}/manifests/{id}/publish",         "Publish manifest, archive previous  →  200  (stub)",                 LAERDAL_RED),
    ("GET",    "/api/manifests/published/{orgCode}",                        "Training app: fetch active manifest snapshot  →  200 / 404  (stub)", BLUE_ACCENT),
]

for i, (method, route, desc, color) in enumerate(endpoints):
    ey = Inches(1.05) + i * Inches(0.88)
    box(s, Inches(0.5), ey, Inches(0.85), Inches(0.62),
        fill_color=color)
    txt(s, method, Inches(0.5), ey, Inches(0.85), Inches(0.62),
        size=Pt(13), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    box(s, Inches(1.45), ey, Inches(11.1), Inches(0.62),
        fill_color=ACCENT, line_color=color, line_width=Pt(1))
    txt(s, route, Inches(1.6), ey + Inches(0.02), Inches(5.5), Inches(0.35),
        size=Pt(13), bold=True, color=WHITE)
    txt(s, desc, Inches(1.6), ey + Inches(0.35), Inches(10.7), Inches(0.25),
        size=Pt(11), color=LIGHT_GREY)

txt(s, "Endpoints marked (stub) are designed but not yet wired up — all infrastructure (repositories, entities) is ready.",
    Inches(0.5), Inches(7.15), Inches(12.5), Inches(0.35),
    size=Pt(12), color=YELLOW, italic=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — Section: Auth
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
section_header(s, "06", "Authentication & Authorization", "OIDC · JWT roles · Service-to-Service")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — Auth design
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s, DARK_BG)
box(s, Inches(0), Inches(0), Inches(0.25), SLIDE_H, fill_color=LAERDAL_RED)

txt(s, "Auth Design", Inches(0.6), Inches(0.2), Inches(10), Inches(0.6),
    size=Pt(32), bold=True, color=WHITE)
divider(s, Inches(0.85), color=LAERDAL_RED)

card(s, Inches(0.5), Inches(1.1), Inches(3.9), Inches(5.9), "React Frontend (OIDC)", [
    "1. User hits POST /api/auth/login",
    "2. Redirect to Laerdal OIDC provider",
    "3. Callback returns JWT access token",
    "",
    "JWT claims:",
    "  sub  → user ID",
    "  org  → primary organization ID",
    "  roles → laerdal_admin | org_admin",
    "          | org_editor",
    "",
    "Middleware validates JWT signature",
    "using OIDC public key (cached + refreshed).",
], accent_color=LAERDAL_RED)

card(s, Inches(4.55), Inches(1.1), Inches(4.2), Inches(5.9), "Role-Based Authorization", [
    "Laerdal Admin (role=laerdal_admin)",
    "  Full CRUD on any organization",
    "  Manage manifests across all orgs",
    "",
    "Hospital Admin (role=org_admin)",
    "  CRUD scoped to own org + descendants",
    "  Manage manifests for their org only",
    "  Cannot modify parent org",
    "",
    "Hospital Editor",
    "  Draft & publish manifests",
    "  Read-only on org settings",
    "",
    "[Authorize(Policy=\"LaerdalAdminOnly\")]",
], accent_color=YELLOW)

card(s, Inches(8.9), Inches(1.1), Inches(4.1), Inches(5.9), "Training App (S2S)", [
    "OAuth 2.0 Client Credentials flow",
    "",
    "1. Training app calls POST /oauth/token",
    "   with client_id + client_secret",
    "2. Receives access token",
    "3. Calls GET /api/manifests/published",
    "   /{orgCode}  with Bearer token",
    "",
    "Scope: manifest:read:{orgCode}",
    "  → token is locked to one org",
    "  → cannot read other orgs' manifests",
    "",
    "Client ID/Secret issued per deployment.",
], accent_color=BLUE_ACCENT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 16 — Section: Persistence & Versioning
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
section_header(s, "07", "Persistence & Versioning", "Schema · Indexes · Learner session safety")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 17 — DB Schema + versioning
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s, DARK_BG)
box(s, Inches(0), Inches(0), Inches(0.25), SLIDE_H, fill_color=LAERDAL_RED)

txt(s, "Database Schema & Versioning Strategy", Inches(0.6), Inches(0.2), Inches(12), Inches(0.6),
    size=Pt(30), bold=True, color=WHITE)
divider(s, Inches(0.82), color=LAERDAL_RED)

# Organizations table
card(s, Inches(0.5), Inches(1.0), Inches(4.5), Inches(5.8), "Organizations Table", [
    "Id             uniqueidentifier  PK",
    "ParentId       uniqueidentifier  NULL  FK→self",
    "Name           nvarchar(255)     NOT NULL",
    "Code           nvarchar(50)      NOT NULL",
    "Type           int               NOT NULL",
    "ExternalId     nvarchar(255)     NULL",
    "IsActive       bit               DEFAULT 1",
    "CreatedAt      datetime2         NOT NULL",
    "UpdatedAt      datetime2         NOT NULL",
    "",
    "INDEX: UQ_Code_ParentId  (Code, ParentId)  UNIQUE",
    "FK:    ON DELETE RESTRICT  (no orphan children)",
], accent_color=LAERDAL_RED, body_size=Pt(12))

# Manifests table
card(s, Inches(5.15), Inches(1.0), Inches(4.5), Inches(5.8), "Manifests Table", [
    "Id             uniqueidentifier  PK",
    "OrganizationId uniqueidentifier  NOT NULL  FK",
    "Name           nvarchar(255)     NOT NULL",
    "Version        nvarchar(20)      NOT NULL",
    "Status         int               NOT NULL",
    "Content        nvarchar(max)     NOT NULL  (JSON)",
    "PublishedAt    datetime2         NULL",
    "CreatedAt      datetime2         NOT NULL",
    "UpdatedAt      datetime2         NOT NULL",
    "",
    "INDEX: UQ_Manifest_OrgVersion  (OrgId, Version)  UNIQUE",
    "INDEX: IX_Manifest_OrgStatus   (OrgId, Status)",
], accent_color=BLUE_ACCENT, body_size=Pt(12))

# Versioning strategy card
card(s, Inches(9.8), Inches(1.0), Inches(3.2), Inches(5.8), "Learner Session Safety", [
    "Problem:",
    "Admin publishes v2 while a learner",
    "is mid-session on v1.",
    "",
    "Solution:",
    "1. Session stores ManifestId (pinned)",
    "2. Publish = INSERT new row",
    "   (never UPDATE published content)",
    "3. Old row → Status = Archived",
    "   (never deleted)",
    "4. Learner reads by Id — always",
    "   gets the same snapshot.",
    "",
    "Published rows are write-once.",
], accent_color=GREEN, body_size=Pt(12))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 18 — What's Next
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s, DARK_BG)
box(s, Inches(0), Inches(0), Inches(0.25), SLIDE_H, fill_color=LAERDAL_RED)

txt(s, "What's Next", Inches(0.6), Inches(0.2), Inches(10), Inches(0.6),
    size=Pt(36), bold=True, color=WHITE)
divider(s, Inches(0.85), color=LAERDAL_RED)

next_items = [
    ("Manifest endpoints", " — wire up create draft, publish, training-app fetch (infrastructure already ready)"),
    ("OIDC middleware",     " — JWT validation, policy-based role enforcement, scope checks"),
    ("PATCH /organizations/{id}", " — UpdateOrganizationCommand + handler (entity method + repo method exist)"),
    ("Input validation",   " — FluentValidation on request models (required fields, code format)"),
    ("Integration tests",  " — test the full Create Org flow against a real test database"),
    ("Audit / events",     " — domain events (OrganizationCreated, ManifestPublished) → audit log or webhooks"),
    ("Pagination",         " — cursor- or offset-based paging on GET /organizations for large datasets"),
    ("NULL index gap",     " — filtered unique index for root org codes (ParentId IS NULL) in SQL Server"),
]

for i, (title, desc) in enumerate(next_items):
    col = 0 if i < 4 else 1
    row = i if i < 4 else i - 4
    ix = Inches(0.5) + col * Inches(6.4)
    iy = Inches(1.15) + row * Inches(1.45)
    box(s, ix, iy, Inches(6.0), Inches(1.2),
        fill_color=ACCENT, line_color=LAERDAL_RED, line_width=Pt(1))
    txt(s, title, ix + Inches(0.15), iy + Inches(0.1), Inches(5.7), Inches(0.42),
        size=Pt(15), bold=True, color=WHITE)
    txt(s, desc.strip(), ix + Inches(0.15), iy + Inches(0.52), Inches(5.7), Inches(0.6),
        size=Pt(12), color=LIGHT_GREY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 19 — Summary
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s, DARK_BG)
box(s, Inches(0), Inches(0), Inches(0.25), SLIDE_H, fill_color=LAERDAL_RED)

txt(s, "Summary", Inches(0.6), Inches(0.3), Inches(10), Inches(0.7),
    size=Pt(40), bold=True, color=WHITE)
divider(s, Inches(1.0), color=LAERDAL_RED)

summary_points = [
    ("Clean Architecture", "4 layers, dependencies point inward only. Domain has zero framework references."),
    ("Domain-first design", "Business rules (immutable manifests, code uniqueness, CanBeDeleted) live on entities — not controllers."),
    ("CQRS + MediatR", "Every operation is an explicit Command or Query. Controllers are thin dispatchers."),
    ("EF Core + migrations", "Fluent configuration in Infrastructure keeps Domain clean. One migration creates the full schema."),
    ("Immutable versioning", "Published manifests are write-once INSERT rows. Learner sessions pin a ManifestId and never see drift."),
    ("Extensible stubs", "Manifest endpoints, OIDC auth, and audit events are designed and documented — ready to build next."),
]

for i, (title, desc) in enumerate(summary_points):
    sy = Inches(1.15) + i * Inches(1.02)
    box(s, Inches(0.5), sy, Inches(0.45), Inches(0.75),
        fill_color=LAERDAL_RED)
    txt(s, str(i + 1), Inches(0.5), sy, Inches(0.45), Inches(0.75),
        size=Pt(18), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, title + ":", Inches(1.1), sy + Inches(0.08), Inches(3.5), Inches(0.4),
        size=Pt(16), bold=True, color=WHITE)
    txt(s, desc, Inches(4.6), sy + Inches(0.08), Inches(8.6), Inches(0.65),
        size=Pt(14), color=LIGHT_GREY)


# ══════════════════════════════════════════════════════════════════════════════
# Save
# ══════════════════════════════════════════════════════════════════════════════
out = "/Users/ulyanahassan/VScodeProjects/LaerdalImplementation/LaerdalImplementation.pptx"
prs.save(out)
print(f"Saved: {out}")
